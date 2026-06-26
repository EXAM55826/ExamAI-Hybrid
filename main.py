import os
import json
import asyncio
import io
import re
from typing import Optional, List
from fastapi import FastAPI, UploadFile, File, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
import httpx
app = FastAPI(title="ExamAI Híbrido API", version="1.0.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
@app.exception_handler(Exception)
async def global_exception_handler(request: Request, exc: Exception):
    return JSONResponse(
        status_code=500,
        content={"error": f"Error interno del servidor: {str(exc)}"}
    )
def extract_text_from_pdf(content: bytes) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            pages_text = []
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages_text.append(text)
            if not pages_text:
                raise ValueError("El PDF solo contiene imágenes escaneadas. Sube un archivo con texto seleccionable.")
            return "\n\n".join(pages_text)
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"No se pudo leer el PDF: {str(e)}")
def extract_text_from_docx(content: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        if not paragraphs:
            raise ValueError("El archivo Word no contiene texto legible.")
        return "\n\n".join(paragraphs)
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"No se pudo leer el archivo Word: {str(e)}")
def extract_text_from_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        slides_text = []
        for slide in prs.slides:
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                slides_text.append("\n".join(slide_texts))
        if not slides_text:
            raise ValueError("El PowerPoint no contiene texto legible.")
        return "\n\n---\n\n".join(slides_text)
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"No se pudo leer el PowerPoint: {str(e)}")
def extract_text_from_txt(content: bytes) -> str:
    try:
        return content.decode("utf-8", errors="replace")
    except Exception as e:
        raise ValueError(f"No se pudo leer el archivo de texto: {str(e)}")
@app.get("/api/healthz")
async def health():
    return {"status": "ok"}
@app.post("/api/upload")
async def upload_file(file: UploadFile = File(...)):
    try:
        await file.seek(0)
        content = await file.read()
        if len(content) == 0:
            raise HTTPException(status_code=400, detail="El archivo está vacío o no se pudo cargar correctamente.")
        filename = file.filename or ""
        ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
        try:
            if ext == "pdf":
                text = extract_text_from_pdf(content)
            elif ext == "docx":
                text = extract_text_from_docx(content)
            elif ext in ("pptx", "ppt"):
                text = extract_text_from_pptx(content)
            elif ext == "txt":
                text = extract_text_from_txt(content)
            else:
                raise HTTPException(
                    status_code=400,
                    detail=f"Formato no soportado: .{ext}. Usa PDF, Word (.docx), PowerPoint (.pptx) o Texto (.txt)."
                )
        except ValueError as e:
            raise HTTPException(status_code=400, detail=str(e))
        word_count = len(text.split())
        return {
            "text": text,
            "word_count": word_count,
            "filename": filename,
            "format": ext.upper()
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error procesando el archivo: {str(e)}")
def build_chunks(text: str, chunk_size: int = 800, overlap: int = 150) -> List[str]:
    words = text.split()
    if len(words) <= chunk_size:
        return [text]
    chunks = []
    start = 0
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        if end >= len(words):
            break
        start = end - overlap
    return chunks

SYSTEM_PROMPT_EXAM = """Eres un generador experto de exámenes académicos en español bajo el esquema de 'Extracción Rígida de Contexto'. 
REGLA ABSOLUTA Y OBLIGATORIA: Solo puedes usar información factual explícita contenida en el documento proporcionado. 
Si un concepto, fecha, dato o afirmación no existe exactamente en el documento, NO puedes inventarlo, asumirlo ni deducirlo bajo ninguna circunstancia. Queda estrictamente PROHIBIDO añadir cualquier conocimiento externo a este contexto.
Responde ÚNICAMENTE con JSON válido, sin texto adicional antes o después."""

SYSTEM_PROMPT_DETECT = """Eres un detector experto de plagio y contenido generado por IA bajo el esquema de 'Extracción Rígida de Contexto'.
Analiza el texto con máximo rigor científico basándote única y exclusivamente en los patrones gramaticales y semánticos del contenido proporcionado, sin añadir suposiciones ni interpretaciones externas.
Responde ÚNICAMENTE con JSON válido, sin texto adicional antes o después."""

async def call_gemini(api_keys: List[str], prompt: str, system_prompt: str, temperature: float = 0.0) -> dict:
    keys = [k.strip() for k in api_keys if k.strip()]
    last_error = None
    for attempt, key in enumerate(keys[:3]):
        try:
            url = f"https://generativelanguage.googleapis.com/v1/models/gemini-1.5-flash:generateContent?key={key}"
            payload = {
                "system_instruction": {"parts": [{"text": system_prompt}]},
                "contents": [{"parts": [{"text": prompt}]}],
                "generationConfig": {
                    "temperature": temperature,
                    "topP": 1,
                    "responseMimeType": "application/json"
                }
            }
            async with httpx.AsyncClient(timeout=90) as client:
                resp = await client.post(url, json=payload)
                if resp.status_code == 429:
                    last_error = "Cuota agotada"
                    continue
                resp.raise_for_status()
                data = resp.json()
                text = data["candidates"][0]["content"]["parts"][0]["text"]
                return json.loads(text)
        except (json.JSONDecodeError, KeyError):
            if "text" in locals():
                raise ValueError(f"La IA devolvió formato inválido: {text[:200]}")
            last_error = "Respuesta inválida"
        except Exception as e:
            last_error = str(e)
    raise HTTPException(status_code=503, detail=f"Gemini no disponible: {last_error}")

async def call_groq(api_keys: List[str], prompt: str, system_prompt: str, temperature: float = 0.0) -> dict:
    keys = [k.strip() for k in api_keys if k.strip()]
    last_error = None
    for key in keys[:3]:
        try:
            url = "https://api.groq.com/openai/v1/chat/completions"
            payload = {
              
                "model": "llama3-8b-8192",
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt}
                ],
                "temperature": temperature,
                "top_p": 1,
                "response_format": {"type": "json_object"}
            }
            headers = {"Authorization": f"Bearer {key}", "Content-Type": "application/json"}
            async with httpx.AsyncClient(timeout=90) as client:
                resp = await client.post(url, json=payload, headers=headers)
                if resp.status_code == 429:
                    last_error = "Cuota agotada"
                    continue
                resp.raise_for_status()
                data = resp.json()
                text = data["choices"][0]["message"]["content"]
                return json.loads(text)
        except (json.JSONDecodeError, KeyError):
            last_error = "Respuesta inválida"
        except Exception as e:
            last_error = str(e)
    raise HTTPException(status_code=503, detail=f"Groq no disponible: {last_error}")

async def call_cohere(api_keys: List[str], prompt: str, system_prompt: str, temperature: float = 0.0) -> dict:
    keys = [k.strip() for k in api_keys if k.strip()]
    last_error = None
    for key in keys[:3]:
        try:
            url = "https://api.cohere.com/v2/chat"
            payload = {
                "model": "command-r-plus-08-2024",
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt}
                ],
                "temperature": temperature,
                "response_format": {"type": "json_object"}
            }
            headers = {"Authorization": f"Bearer {key}", "Content-Type": "application/json"}
            async with httpx.AsyncClient(timeout=90) as client:
                resp = await client.post(url, json=payload, headers=headers)
                if resp.status_code == 429:
                    last_error = "Cuota agotada"
                    continue
                resp.raise_for_status()
                data = resp.json()
                text = data["message"]["content"][0]["text"]
                return json.loads(text)
        except (json.JSONDecodeError, KeyError):
            last_error = "Respuesta inválida"
        except Exception as e:
            last_error = str(e)
    raise HTTPException(status_code=503, detail=f"Cohere no disponible: {last_error}")

async def call_ai(provider: str, api_keys: List[str], prompt: str, system_prompt: str, temperature: float = 0.0) -> dict:
    if provider == "gemini":
        return await call_gemini(api_keys, prompt, system_prompt, temperature)
    elif provider == "groq":
        return await call_groq(api_keys, prompt, system_prompt, temperature)
    elif provider == "cohere":
        return await call_cohere(api_keys, prompt, system_prompt, temperature)
    else:
        raise HTTPException(status_code=400, detail=f"Proveedor no soportado: {provider}")
@app.post("/api/generate-exam")
async def generate_exam(request: Request):
    try:
        body = await request.json()
        text = body.get("text", "")
        provider = body.get("provider", "gemini")
        api_keys = body.get("api_keys", [])
        difficulty = body.get("difficulty", "media")
        subject = body.get("subject", "")
        question_count = int(body.get("question_count", 40))
        variation_seed = body.get("variation_seed", 0)
        question_types = body.get("question_types", ["multiple", "truefalse", "complete", "match"])
        rotate_providers = bool(body.get("rotate_providers", False))
        api_keys_dict = body.get("api_keys_dict", {})
        active_providers = []
        if rotate_providers and isinstance(api_keys_dict, dict):
            for p in ["gemini", "groq", "cohere"]:
                keys = api_keys_dict.get(p, [])
                if isinstance(keys, list):
                    filtered_keys = [k.strip() for k in keys if k.strip()]
                    if filtered_keys:
                        active_providers.append((p, filtered_keys))
        if not text.strip():
            raise HTTPException(status_code=400, detail="No hay texto para analizar.")
       
        if not api_keys and not active_providers:
            raise HTTPException(status_code=400, detail="Se requiere al menos una API key o proveedores configurados para rotación.")
        chunks = build_chunks(text, 800, 150)
        context = chunks[0] if len(chunks) == 1 else "\n\n[...]\n\n".join(chunks[:3])
        variation_hint = ""
        if variation_seed > 0:
            variation_hint = f"\n\nIMPORTANTE: Esta es la variación #{variation_seed}. Genera un conjunto COMPLETAMENTE DIFERENTE de preguntas, explorando otros temas y conceptos del documento que no hayas cubierto anteriormente."
        all_questions = []
        batch_size = 10
        batches = (question_count + batch_size - 1) // batch_size
        for batch_idx in range(batches):
            batch_num = batch_idx + 1
            remaining = question_count - len(all_questions)
            current_batch = min(batch_size, remaining)
            used_topics = json.dumps([q.get("question", "")[:60] for q in all_questions[:5]]) if all_questions else "[]"
            type_label_map = {
                "multiple":  "Opción múltiple (A/B/C/D)",
                "truefalse": "Verdadero o Falso",
                "complete":  "Completar espacio en blanco",
                "match":     "Relacionar conceptos (pares)",
            }
            valid_types = [t for t in question_types if t in type_label_map]
            if not valid_types:
                valid_types = ["multiple"]
            selected_str = " | ".join(type_label_map[t] for t in valid_types)
            type_rule = (
                f"TIPOS PERMITIDOS — Genera ÚNICAMENTE: {selected_str}. "
                f"PROHIBIDO cualquier otro tipo. Distribuye equitativamente entre los tipos seleccionados."
            )
            fmt_parts = []
            base_id = len(all_questions) + 1
            if "multiple" in valid_types:
                fmt_parts.append(
                    f'{{"id": {base_id}, "type": "multiple", "question": "¿Cuál es...?", '
                    f'"options": ["A) opción1", "B) opción2", "C) opción3", "D) opción4"], "answer": "A", "explanation": "..."}}'
                )
            if "truefalse" in valid_types:
                fmt_parts.append(
                    f'{{"id": {base_id}, "type": "truefalse", "question": "Afirmación sobre el tema.", '
                    f'"options": ["A) Verdadero", "B) Falso"], "answer": "B", "explanation": "..."}}'
                )
            if "complete" in valid_types:
                fmt_parts.append(
                    f'{{"id": {base_id}, "type": "complete", "question": "El proceso de _____ permite...", '
                    f'"options": [], "answer": "fotosíntesis", "explanation": "..."}}'
                )
            if "match" in valid_types:
                fmt_parts.append(
                    f'{{"id": {base_id}, "type": "match", "question": "Une cada concepto con su definición:", '
                    f'"options": [], "answer": "Ver pares", '
                    f'"pairs": [{{"left": "Concepto A", "right": "Def A"}}, {{"left": "Concepto B", "right": "Def B"}}, '
                    f'{{"left": "Concepto C", "right": "Def C"}}, {{"left": "Concepto D", "right": "Def D"}}], '
                    f'"explanation": "..."}}'
                )
            fmt_examples = ",\n    ".join(fmt_parts)
            prompt = f"""Documento académico:
\"\"\"
{context}
\"\"\"
{variation_hint}
Genera exactamente {current_batch} preguntas de examen (lote {batch_num} de {batches}).
Dificultad: {difficulty}
{f'Área: {subject}' if subject else ''}
Temas ya cubiertos (evitar repetir): {used_topics}
{type_rule}
FORMATO JSON ESTRICTO — usa el tipo correcto para cada pregunta:
{{
  "questions": [
    {fmt_examples}
  ]
}}
REGLAS:
- Solo usa información del documento
- El campo "type" es obligatorio en cada pregunta
- multiple: 4 opciones (A/B/C/D), "answer" es la letra
- truefalse: 2 opciones (A) Verdadero, B) Falso), "answer" es la letra
- complete: "options" vacío [], "answer" es la palabra/frase exacta del espacio
- match: "options" vacío [], "pairs" con exactamente 4 pares {{left, right}}, "answer" es "Ver pares"
- No repitas preguntas del lote anterior"""
           
            if rotate_providers and active_providers:
                current_provider, current_keys = active_providers[batch_idx % len(active_providers)]
            else:
                current_provider, current_keys = provider, api_keys
            result = await call_ai(current_provider, current_keys, prompt, SYSTEM_PROMPT_EXAM, temperature=0.0)
            batch_questions = result.get("questions", [])
            all_questions.extend(batch_questions)
            if len(all_questions) >= question_count:
                break
        return {
            "questions": all_questions[:question_count],
            "total": len(all_questions[:question_count]),
            "provider": provider
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generando examen: {str(e)}")
@app.post("/api/detect")
async def detect_plagiarism(request: Request):
    try:
        body = await request.json()
        text = body.get("text", "")
        provider = body.get("provider", "gemini")
        api_keys = body.get("api_keys", [])
        if not text.strip():
            raise HTTPException(status_code=400, detail="No hay texto para analizar.")
        if not api_keys:
            raise HTTPException(status_code=400, detail="Se requiere al menos una API key.")
        chunks = build_chunks(text, 800, 150)
        context = chunks[0]
        prompt = f"""Analiza el siguiente texto académico para detectar:
1. Contenido generado por IA (probabilidad y señales)
2. Posible plagio (patrones de escritura inusuales, inconsistencias de estilo)
3. Métricas de perplejidad y burstiness
4. Análisis de originalidad
5. Fragmentos específicos de texto (oraciones o párrafos) que muestren fuerte evidencia de ser IA, plagio u originales.
Texto a analizar:
\"\"\"
{context}
\"\"\"
Responde con este JSON EXACTO:
{{
  "ai_probability": 0.75,
  "plagiarism_probability": 0.30,
  "originality_score": 0.85,
  "perplexity_score": 45.2,
  "burstiness_score": 0.62,
  "verdict": "PROBABLE_IA" | "PROBABLE_PLAGIO" | "ORIGINAL" | "SOSPECHOSO",
  "confidence": "alta" | "media" | "baja",
  "ai_indicators": ["lista de señales de IA detectadas"],
  "plagiarism_indicators": ["lista de señales de plagio"],
  "style_analysis": {{
    "sentence_variety": "alta|media|baja",
    "vocabulary_richness": "alta|media|baja",
    "coherence": "alta|media|baja",
    "formality": "formal|semiformal|informal"
  }},
  
  "summary": "Resumen ejecutivo del análisis en 2-3 oraciones",
  "highlighted_segments": [
    {{
      "text": "frase o fragmento exacto que aparece en el texto analizado",
      "type": "ai" | "plagiarism" | "original",
      "probability": 0.85,
      "reason": "explicación de la señal o patrón detectado en este fragmento"
    }}
  ]
}}"""
        result = await call_ai(provider, api_keys, prompt, SYSTEM_PROMPT_DETECT, temperature=0.0)
        # Asegurar que existan los campos clave
        if "highlighted_segments" not in result:
            result["highlighted_segments"] = []
        return result
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error en detección: {str(e)}")
